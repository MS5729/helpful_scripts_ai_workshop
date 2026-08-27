from __future__ import annotations

import io
import json
import os
from pathlib import Path
from typing import Any

from dotenv import load_dotenv

load_dotenv()


def parse_file(path: str) -> list[dict[str, Any]]:
    file_path = Path(path)
    suffix = file_path.suffix.lower()
    if suffix == ".txt":
        return [{"text": file_path.read_text(encoding="utf-8"), "page": None}]
    if suffix == ".csv":
        import pandas as pd
        frame = pd.read_csv(file_path)
        return [{"text": " | ".join(f"{k}: {v}" for k, v in row.items() if str(v) != "nan"), "row": i} for i, row in enumerate(frame.to_dict("records"), 1)]
    if suffix in {".xlsx", ".xls"}:
        import pandas as pd
        frame = pd.read_excel(file_path)
        return [{"text": " | ".join(f"{k}: {v}" for k, v in row.items() if str(v) != "nan"), "row": i} for i, row in enumerate(frame.to_dict("records"), 1)]
    if suffix == ".pdf":
        import pdfplumber
        with pdfplumber.open(file_path) as pdf:
            return [{"text": page.extract_text() or "", "page": i} for i, page in enumerate(pdf.pages, 1) if page.extract_text()]
    if suffix == ".docx":
        from docx import Document
        document = Document(file_path)
        return [{"text": p.text, "paragraph": i} for i, p in enumerate(document.paragraphs, 1) if p.text.strip()]
    if suffix == ".pptx":
        from pptx import Presentation
        presentation = Presentation(file_path)
        return [{"text": shape.text, "slide": i} for i, slide in enumerate(presentation.slides, 1) for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
    raise ValueError(f"Unsupported file type: {suffix or 'unknown'}")


def chunk_records(records: list[dict[str, Any]], size: int = 800, overlap: int = 100) -> list[dict[str, Any]]:
    if size <= overlap:
        raise ValueError("size must be greater than overlap")
    chunks = []
    step = size - overlap
    for record_index, record in enumerate(records):
        words = record["text"].split()
        for start in range(0, len(words), step):
            text = " ".join(words[start:start + size])
            if text:
                chunks.append({"text": text, "metadata": {"record": record_index, **{k: v for k, v in record.items() if k != "text"}}})
            if start + size >= len(words):
                break
    return chunks


def normalize_row(row: dict[str, Any], *, supplier_id: str, source_system: str, filename: str) -> dict[str, Any]:
    """Map common quality-history columns into a reusable event record."""
    normalized = {"".join(character for character in str(key).lower() if character.isalnum()): value for key, value in row.items()}

    def value(*names: str) -> str:
        for name in names:
            candidate = normalized.get("".join(character for character in name.lower() if character.isalnum()))
            if candidate is not None and str(candidate).strip() and str(candidate).lower() != "nan":
                return str(candidate).strip()
        return ""

    severity = value("severity", "priority", "risk").lower()
    if severity not in {"low", "medium", "high", "critical"}:
        try:
            score = int(float(severity))
            severity = "critical" if score >= 9 else "high" if score >= 7 else "low" if score <= 3 else "medium"
        except ValueError:
            severity = "medium"
    title = value("title", "failure mode", "issue", "summary", "name")
    description = value("description", "details", "failure description", "text") or title
    return {
        "supplier_id": supplier_id,
        "title": (title or description[:300] or "Imported quality event")[:300],
        "description": description,
        "category": (value("category", "part category", "problem group", "function") or "Uncategorized")[:120],
        "source_system": source_system[:120],
        "source_reference": value("source reference", "reference", "jira", "id") or None,
        "severity": severity,
        "status": value("status", "state") or "open",
        "source_filename": filename,
    }


def output_json(name: str, value: Any) -> None:
    output_dir = Path("workshop_output")
    output_dir.mkdir(exist_ok=True)
    target = output_dir / name
    target.write_text(json.dumps(value, indent=2, default=str), encoding="utf-8")
    print(f"Wrote {target}")


def required_env(name: str) -> str:
    value = os.getenv(name)
    if not value:
        raise RuntimeError(f"Set {name} in .env before running this script")
    return value
